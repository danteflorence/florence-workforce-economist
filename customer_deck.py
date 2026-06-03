"""
Customer-facing presentation generator.

Produces an 8-slide editable PowerPoint (.pptx) per health system, mirroring
the Florence + AMN + Kaiser Permanente reference deck. Brand-aligned with the
Streamlit app (Florence Teal #0ABAB5, Navy #101828, Playfair Display serif headlines,
Inter body sans).

Slide structure:
  1. Cover           — System name + partnership lockup line
  2. Same hours      — TODAY vs WITH FLORENCE per-hour comparison
  3. Conversion      — "Full-time conversion, not agency replacement"
  4. Top facilities  — top 5 facilities by 24-mo savings
  5. ROI summary     — savings ratio + payback
  6. Implementation  — timeline + payback
  7. What this means — closing pitch
  8. Compliance      — sources + disclosure

Also generates a print-ready PDF rendering of the same content (via LibreOffice
fallback if available, otherwise from the HTML exec summary path we already have).

Wired into app.py as a 4th download button next to Excel / PDF / Bundle.
"""
from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Brand palette
TEAL = RGBColor(0x0B, 0xC5, 0xA0)
TEAL_DARK = RGBColor(0x08, 0x94, 0x78)
NAVY = RGBColor(0x0F, 0x1B, 0x2D)
NAVY_SOFT = RGBColor(0x1A, 0x2A, 0x44)
GRAY = RGBColor(0xF4, 0xF6, 0xF8)
GRAY_BORDER = RGBColor(0xE5, 0xE8, 0xEE)
INK = RGBColor(0x0F, 0x1B, 0x2D)
MUTED = RGBColor(0x5B, 0x66, 0x75)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Standard 16:9 slide is 13.333" x 7.5"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


@dataclass
class DeckInputs:
    system_name: str
    n_facilities: int
    n_states: int
    rn_need: float
    agency_hourly: float                # median agency premium / hour ($)
    florence_hourly_net: float          # Florence effective net / hour ($)
    annual_savings: float               # net savings per year ($)
    term_savings: float                 # net savings over 24 mo ($)
    term_florence_fee: float            # Florence 24-mo fee total ($)
    savings_ratio: float                # term_savings / term_florence_fee
    top_facilities: pd.DataFrame        # top 5 facilities by savings
    target_offset_pct: float            # e.g. 0.50


def _fmt_big(v: float) -> str:
    """Money with two decimals at the right magnitude, so small values never
    collapse to a misleading '$0M'."""
    v = float(v or 0)
    if abs(v) >= 1e9:
        return f"${v/1e9:,.2f}B"
    if abs(v) >= 1e6:
        return f"${v/1e6:,.2f}M"
    if abs(v) >= 1e3:
        return f"${v/1e3:,.2f}K"
    return f"${v:,.2f}"


def _add_text(slide, x, y, w, h, text, *,
              font="Inter", size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    """Add a text box with simple formatting."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    return tb


def _add_rect(slide, x, y, w, h, *, fill=GRAY, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.shadow.inherit = False
    return sh


def _eyebrow(slide, x, y, text):
    return _add_text(slide, x, y, Inches(8), Inches(0.3),
                     text.upper(), font="Inter", size=10, bold=True,
                     color=TEAL_DARK)


def _h1(slide, x, y, text, w_in=10):
    """Big editorial serif headline."""
    return _add_text(slide, x, y, Inches(w_in), Inches(1.6),
                     text, font="Playfair Display", size=44, bold=True, color=NAVY)


def _h3(slide, x, y, text, w_in=10):
    return _add_text(slide, x, y, Inches(w_in), Inches(0.7),
                     text, font="Playfair Display", size=22, bold=True, color=NAVY)


def _body(slide, x, y, text, w_in=10, color=MUTED, size=14):
    return _add_text(slide, x, y, Inches(w_in), Inches(2),
                     text, font="Inter", size=size, color=color)


def _brand_strip(slide):
    """Florence brand mark + section tag, top of every slide."""
    # Teal F box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(0.3), Inches(0.35), Inches(0.35),
    )
    box.fill.solid(); box.fill.fore_color.rgb = TEAL
    box.line.fill.background()
    box.shadow.inherit = False
    box.adjustments[0] = 0.2
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(0.35), Inches(0.4))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Pt(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "F"
    r.font.name = "Inter"; r.font.size = Pt(16); r.font.bold = True
    r.font.color.rgb = WHITE
    # Wordmark
    _add_text(slide, Inches(0.95), Inches(0.32), Inches(2), Inches(0.4),
              "Florence", font="Playfair Display", size=16, bold=True, color=NAVY)


def _section_tag(slide, text):
    _add_text(slide, Inches(8.5), Inches(0.35), Inches(4.5), Inches(0.3),
              text, font="Inter", size=9, bold=True, color=MUTED,
              align=PP_ALIGN.RIGHT)


def _slide_cover(prs: Presentation, d: DeckInputs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _brand_strip(s)
    _section_tag(s, "CONFIDENTIAL PROPOSAL  ·  FLORENCE WORKFORCE ECONOMIST")
    _eyebrow(s, Inches(0.5), Inches(2.3), f"For {d.system_name}")
    _h1(s, Inches(0.5), Inches(2.7),
        "A Pathway to Permanent\nRN Capacity.", w_in=12)
    _body(s, Inches(0.5), Inches(5.7),
          f"Florence delivers {d.rn_need:,.0f} permanent international RNs to "
          f"{d.system_name} across {d.n_facilities} facilit{'ies' if d.n_facilities != 1 else 'y'} in "
          f"{d.n_states} state{'s' if d.n_states != 1 else ''}. "
          f"Same hours. Different price. Different outcome.",
          w_in=10, size=14)
    _add_text(s, Inches(0.5), Inches(7.0), Inches(4), Inches(0.3),
              "florenceedu.com", font="Inter", size=10, color=TEAL_DARK)


def _slide_same_hours(prs, d: DeckInputs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _brand_strip(s); _section_tag(s, f"01  ·  THE OPPORTUNITY  ·  FOR {d.system_name.upper()}")
    _eyebrow(s, Inches(0.5), Inches(1.0), "01 · The opportunity")
    _h1(s, Inches(0.5), Inches(1.4), "Same hours. Two prices.")
    _body(s, Inches(0.5), Inches(2.7),
          f"The agency RN hours {d.system_name} used last year will recur. "
          f"The choice is what to pay per hour, and what to get for it.",
          w_in=11, size=14)

    # Two cards
    y = Inches(3.7); h = Inches(2.6)
    # TODAY (gray)
    today = _add_rect(s, Inches(0.5), y, Inches(5.6), h, fill=GRAY, line=GRAY_BORDER)
    _add_text(s, Inches(0.8), Inches(3.85), Inches(3), Inches(0.3),
              "TODAY", font="Inter", size=9, bold=True, color=MUTED)
    _add_text(s, Inches(0.8), Inches(4.15), Inches(5), Inches(1.1),
              f"${d.agency_hourly:,.2f}", font="Playfair Display", size=48, bold=True, color=NAVY)
    _add_text(s, Inches(3.9), Inches(4.55), Inches(2), Inches(0.5),
              "/hour", font="Playfair Display", size=20, color=MUTED)
    _add_text(s, Inches(0.8), Inches(5.2), Inches(5), Inches(0.4),
              "Agency staffing premium.", font="Playfair Display", size=14, bold=True, color=NAVY)
    _add_text(s, Inches(0.8), Inches(5.55), Inches(5), Inches(1),
              "Contingent labor. No continuity of unit, panel, or workforce planning. "
              "Premium recurs every cycle.",
              font="Inter", size=10, color=MUTED)

    # Arrow
    _add_text(s, Inches(6.2), Inches(4.7), Inches(0.8), Inches(0.6),
              "→", font="Playfair Display", size=32, color=TEAL, align=PP_ALIGN.CENTER)

    # WITH FLORENCE (teal)
    with_f = _add_rect(s, Inches(7.0), y, Inches(5.8), h, fill=TEAL)
    _add_text(s, Inches(7.3), Inches(3.85), Inches(3), Inches(0.3),
              "WITH FLORENCE", font="Inter", size=9, bold=True, color=WHITE)
    _add_text(s, Inches(7.3), Inches(4.15), Inches(5), Inches(1.1),
              f"${d.florence_hourly_net:,.2f}", font="Playfair Display", size=48, bold=True, color=WHITE)
    _add_text(s, Inches(10.4), Inches(4.55), Inches(2), Inches(0.5),
              "/hour", font="Playfair Display", size=20, color=WHITE)
    _add_text(s, Inches(7.3), Inches(5.2), Inches(5), Inches(0.4),
              "Permanent RN capacity.", font="Playfair Display", size=14, bold=True, color=WHITE)
    _add_text(s, Inches(7.3), Inches(5.55), Inches(5), Inches(1),
              f"Full-time {d.system_name} employees on multi-year contracts. "
              "Fee payable on successful employment start; replacement protection for early attrition.",
              font="Inter", size=10, color=WHITE)

    # Navy banner — annual savings
    banner = _add_rect(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), fill=NAVY)
    _add_text(s, Inches(0.8), Inches(6.65), Inches(7), Inches(0.4),
              "Annual savings opportunity · net of Florence fees and FICA equivalence",
              font="Inter", size=11, color=WHITE)
    _add_text(s, Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.5),
              _fmt_big(d.annual_savings),
              font="Playfair Display", size=22, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)


def _slide_conversion(prs, d: DeckInputs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _brand_strip(s); _section_tag(s, f"02  ·  THE STRUCTURAL DIFFERENCE")
    _eyebrow(s, Inches(0.5), Inches(1.0), "02 · The structural difference")
    _h1(s, Inches(0.5), Inches(1.4), "This is full-time conversion\n— not agency replacement.")

    # Two-pane comparison
    y = Inches(3.8); h = Inches(3.2)
    pane_w = Inches(5.8)

    # LEFT — Agency Model
    _add_rect(s, Inches(0.5), y, pane_w, h, fill=GRAY, line=GRAY_BORDER)
    _add_text(s, Inches(0.85), Inches(3.95), Inches(4), Inches(0.3),
              "THE AGENCY MODEL TODAY", font="Inter", size=9, bold=True, color=MUTED)
    _add_text(s, Inches(0.85), Inches(4.25), Inches(5), Inches(0.5),
              "Contingent.\nRecycled.\nPremium-priced.",
              font="Playfair Display", size=18, bold=True, color=NAVY)
    bullets1 = [
        "• No continuity of unit, team, or panel",
        f"• ${d.agency_hourly:.2f}/hr staffing premium on every RN hour",
        "• Outside the labor-partnership framework",
        "• Cycle repeats every fiscal year",
    ]
    _add_text(s, Inches(0.85), Inches(5.7), Inches(5), Inches(1.5),
              "\n".join(bullets1), font="Inter", size=11, color=MUTED)

    # RIGHT — Florence Pathway (teal)
    _add_rect(s, Inches(7.0), y, pane_w, h, fill=TEAL)
    _add_text(s, Inches(7.35), Inches(3.95), Inches(4), Inches(0.3),
              "THE FLORENCE PATHWAY", font="Inter", size=9, bold=True, color=WHITE)
    _add_text(s, Inches(7.35), Inches(4.25), Inches(5), Inches(0.5),
              "Permanent.\nAligned.\nRepeatable.",
              font="Playfair Display", size=18, bold=True, color=WHITE)
    bullets2 = [
        f"• {d.system_name}'s standard hiring & onboarding",
        f"• Full-time {d.system_name} employee, {d.system_name} comp & benefits",
        "• 3-year minimum tenure",
        "• Labor-partnership aligned",
    ]
    _add_text(s, Inches(7.35), Inches(5.7), Inches(5), Inches(1.5),
              "\n".join(bullets2), font="Inter", size=11, color=WHITE)


def _slide_top_facilities(prs, d: DeckInputs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _brand_strip(s); _section_tag(s, f"03  ·  PER-FACILITY DETAIL")
    _eyebrow(s, Inches(0.5), Inches(1.0), "03 · Per-facility detail")
    _h1(s, Inches(0.5), Inches(1.4), "Top facilities by 24-month savings.")

    # Header row
    cols_y = Inches(3.3)
    headers = [
        ("Facility", Inches(0.5), Inches(5.5)),
        ("State", Inches(6.1), Inches(0.8)),
        ("RNs", Inches(7.0), Inches(0.8)),
        ("24-mo savings", Inches(8.0), Inches(2.0)),
        ("Fee /RN/mo", Inches(10.1), Inches(1.6)),
        ("ROI", Inches(11.8), Inches(1.0)),
    ]
    for label, x, w in headers:
        _add_text(s, x, cols_y, w, Inches(0.3), label,
                  font="Inter", size=10, bold=True, color=MUTED)

    # Divider line
    _add_rect(s, Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.015), fill=GRAY_BORDER)

    # Rows
    top5 = d.top_facilities.head(5)
    for i, (_, r) in enumerate(top5.iterrows()):
        ry = Inches(3.85 + i * 0.65)
        _add_text(s, Inches(0.5), ry, Inches(5.5), Inches(0.35),
                  str(r.get("name", "—"))[:60],
                  font="Inter", size=12, bold=True, color=NAVY)
        _add_text(s, Inches(0.5), ry + Inches(0.32), Inches(5.5), Inches(0.25),
                  str(r.get("city", "")),
                  font="Inter", size=10, color=MUTED)
        _add_text(s, Inches(6.1), ry, Inches(0.8), Inches(0.35),
                  str(r.get("state", "")), font="Inter", size=12, color=NAVY)
        _add_text(s, Inches(7.0), ry, Inches(0.8), Inches(0.35),
                  f"{float(r.get('rn_need', 0)):.0f}",
                  font="Inter", size=12, color=NAVY)
        _add_text(s, Inches(8.0), ry, Inches(2.0), Inches(0.35),
                  _fmt_big(float(r.get("target_term_net_savings_account", 0))),
                  font="Playfair Display", size=14, bold=True, color=TEAL_DARK)
        _rn = max(float(r.get("rn_need", 0) or 0), 1.0)
        _mo_fee_acct = float(r.get("target_monthly_florence_fee_account", 0) or 0) \
            or (float(r.get("target_term_florence_fee_account", 0) or 0) / 24.0)
        _add_text(s, Inches(10.1), ry, Inches(1.6), Inches(0.35),
                  f"${_mo_fee_acct / _rn:,.2f}",
                  font="Inter", size=12, color=NAVY)
        _add_text(s, Inches(11.8), ry, Inches(1.0), Inches(0.35),
                  f"{float(r.get('target_savings_ratio', 0)):.1f}×",
                  font="Inter", size=12, bold=True, color=TEAL_DARK)


def _slide_roi(prs, d: DeckInputs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _brand_strip(s); _section_tag(s, "04  ·  THE MATH")
    _eyebrow(s, Inches(0.5), Inches(1.0), "04 · The math")
    _h1(s, Inches(0.5), Inches(1.4),
        f"{_fmt_big(d.term_savings)} in net savings over 24 months.")
    _body(s, Inches(0.5), Inches(2.8),
          f"On a Florence fee of {_fmt_big(d.term_florence_fee)} over the term, "
          f"{d.system_name} saves {_fmt_big(d.term_savings)} net — "
          f"a {d.savings_ratio:.1f}× return on every dollar invested with Florence.",
          w_in=12, size=14)

    # Three stat panels
    panels = [
        ("Permanent RNs", f"{d.rn_need:,.0f}", "delivered across the system"),
        ("Annual savings", _fmt_big(d.annual_savings), "to your bottom line"),
        ("Return on investment", f"{d.savings_ratio:.1f}×", "savings to Florence fee"),
    ]
    for i, (label, val, suffix) in enumerate(panels):
        x = Inches(0.5 + i * 4.3)
        _add_rect(s, x, Inches(4.2), Inches(4.0), Inches(2.4), fill=GRAY, line=GRAY_BORDER)
        _add_text(s, x + Inches(0.3), Inches(4.4), Inches(3.5), Inches(0.3),
                  label.upper(), font="Inter", size=9, bold=True, color=MUTED)
        _add_text(s, x + Inches(0.3), Inches(4.8), Inches(3.5), Inches(1),
                  val, font="Playfair Display", size=36, bold=True, color=NAVY)
        _add_text(s, x + Inches(0.3), Inches(5.9), Inches(3.5), Inches(0.4),
                  suffix, font="Inter", size=11, color=MUTED)


def _slide_implementation(prs, d: DeckInputs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _brand_strip(s); _section_tag(s, "05  ·  IMPLEMENTATION")
    _eyebrow(s, Inches(0.5), Inches(1.0), "05 · Implementation")
    _h1(s, Inches(0.5), Inches(1.4), "From signature to first start: 9–12 months.")
    _body(s, Inches(0.5), Inches(2.7),
          "Florence's pipeline runs international RN production end-to-end — "
          "exam prep, visa support, higher education, and bedside practice. "
          f"Onboarding into {d.system_name}'s standard hiring flow.",
          w_in=12, size=13)

    # Timeline cards
    steps = [
        ("Months 0–3", "Cohort selection", "Florence pipeline + system needs assessment"),
        ("Months 3–6", "Visa & licensing", "F-1 transition, NCLEX, state licensure"),
        ("Months 6–9", "Onboarding", "Standard system hiring, comp & benefits"),
        ("Month 9+", "First starts", "Permanent RN production"),
    ]
    for i, (when, what, detail) in enumerate(steps):
        x = Inches(0.5 + i * 3.2)
        _add_rect(s, x, Inches(4.0), Inches(3.0), Inches(2.6), fill=GRAY, line=GRAY_BORDER)
        _add_text(s, x + Inches(0.25), Inches(4.2), Inches(2.5), Inches(0.3),
                  when.upper(), font="Inter", size=9, bold=True, color=TEAL_DARK)
        _add_text(s, x + Inches(0.25), Inches(4.6), Inches(2.5), Inches(0.5),
                  what, font="Playfair Display", size=16, bold=True, color=NAVY)
        _add_text(s, x + Inches(0.25), Inches(5.4), Inches(2.5), Inches(1.1),
                  detail, font="Inter", size=10, color=MUTED)


def _slide_closing(prs, d: DeckInputs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _brand_strip(s); _section_tag(s, f"06  ·  WHAT THIS MEANS FOR {d.system_name.upper()}")
    _eyebrow(s, Inches(0.5), Inches(1.0), f"06 · What this means for {d.system_name}")
    _h1(s, Inches(0.5), Inches(1.4),
        "Permanent RN capacity.\nFICA-aligned economics.\nNo agency dependency.")

    # Closing pitch box
    _add_rect(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), fill=TEAL)
    _add_text(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.4),
              "THE WAY FORWARD", font="Inter", size=10, bold=True, color=WHITE)
    _add_text(s, Inches(0.9), Inches(5.1), Inches(11.5), Inches(1.4),
              f"Florence delivers {d.rn_need:,.0f} permanent RNs to {d.system_name} "
              f"across {d.n_facilities} facilities. Same hours. Different price. "
              f"Different outcome — {_fmt_big(d.term_savings)} in net savings "
              f"on a {_fmt_big(d.term_florence_fee)} investment.",
              font="Playfair Display", size=20, color=WHITE)


def _slide_compliance(prs, d: DeckInputs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _brand_strip(s); _section_tag(s, "07  ·  METHODOLOGY & DISCLOSURES")
    _eyebrow(s, Inches(0.5), Inches(1.0), "07 · Methodology & disclosures")
    _h1(s, Inches(0.5), Inches(1.4), "Sources & assumptions.")

    body = (
        "Universe   ·   CMS Hospital General Information + NASHP Hospital Cost Tool (2020–2024).\n\n"
        "Wage data   ·   BLS OEWS May 2024 MSA-level RN wages (SOC 29-1141) with HCRIS per-hospital and state-level fallback.\n\n"
        "Agency rates   ·   CMS HCRIS NMRC Worksheet S-3 Part II line 01100 (Direct Patient Care contract labor hourly).\n\n"
        f"Pricing   ·   Florence Workforce Restoration Economics v2 — FICA-offset target at {d.target_offset_pct:.0%}. "
        "Florence fee sized as FICA savings ÷ target offset percentage, clamped to $750–$2,000 / RN / month.\n\n"
        "Channel & partner margin   ·   Florence's core rate is protected regardless of channel. "
        "When sold through AMN Healthcare or other distribution partners, the partner's margin is added "
        "on top of Florence's core rate (not deducted from it). Florence collects the full core rate; "
        "the partner collects its added margin. The customer's total price = Florence core rate + partner margin.\n\n"
        "Cohort   ·   F-1 student exemption applies during the nonresident-alien period (IRC §3121(b)(19), IRS Pub 519). "
        "Eligibility must be confirmed by payroll, tax counsel, and immigration counsel.\n\n"
        "FICA-offset reporting   ·   Estimated employer-side FICA offset is calculated on the taxable wage portion of "
        "RN compensation only (not the benefits-loaded HCRIS total comp rate). Eligibility must be validated by payroll, "
        "tax counsel, and immigration counsel. Employee-side FICA benefit is shown separately as a nurse take-home "
        "benefit and is not included in employer ROI unless presented as combined economic value."
    )
    _add_text(s, Inches(0.5), Inches(2.9), Inches(12), Inches(4),
              body, font="Inter", size=11, color=INK)


def build_deck(d: DeckInputs) -> BytesIO:
    """Build the 8-slide deck and return as in-memory BytesIO."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_cover(prs, d)
    _slide_same_hours(prs, d)
    _slide_conversion(prs, d)
    _slide_top_facilities(prs, d)
    _slide_roi(prs, d)
    _slide_implementation(prs, d)
    _slide_closing(prs, d)
    _slide_compliance(prs, d)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def build_deck_from_system_recs(
    system_recs: pd.DataFrame,
    system_name: str,
    target_offset_pct: float = 0.50,
) -> BytesIO:
    """Convenience wrapper: build the deck directly from a slice of recommendations.parquet."""
    if len(system_recs) == 0:
        raise ValueError("No facilities for system " + system_name)
    n_facilities = len(system_recs)
    n_states = system_recs["state"].nunique()
    rn_need = system_recs["rn_need"].sum()
    median_agency_premium = float(system_recs["signal_agency_premium"].median())
    median_florence_hourly = float(system_recs["target_hourly_fee"].median())
    median_fica_per_hour = float(system_recs["target_fica_savings_per_rn_per_month"].median()) / 156
    florence_hourly_net = max(median_florence_hourly - median_fica_per_hour, 0.01)
    term_savings = float(system_recs["target_term_net_savings_account"].sum())
    term_florence_fee = float(system_recs["target_term_florence_fee_account"].sum())
    annual_savings = term_savings / 2
    savings_ratio = term_savings / term_florence_fee if term_florence_fee > 0 else 0
    top_facilities = system_recs.sort_values(
        "target_term_net_savings_account", ascending=False
    )

    d = DeckInputs(
        system_name=system_name,
        n_facilities=n_facilities,
        n_states=n_states,
        rn_need=rn_need,
        agency_hourly=median_agency_premium,
        florence_hourly_net=florence_hourly_net,
        annual_savings=annual_savings,
        term_savings=term_savings,
        term_florence_fee=term_florence_fee,
        savings_ratio=savings_ratio,
        top_facilities=top_facilities,
        target_offset_pct=target_offset_pct,
    )
    return build_deck(d)


if __name__ == "__main__":
    # CLI smoke test: build a deck for Kaiser Permanente
    DATA = Path(__file__).parent / "data"
    recs = pd.read_parquet(DATA / "recommendations.parquet")
    feas = recs[recs["feasible"]].copy()
    kp = feas[feas["health_system_id"] == "kaiser_permanente"]
    print(f"Building deck for Kaiser Permanente ({len(kp)} facilities)…")
    buf = build_deck_from_system_recs(kp, "Kaiser Permanente")
    out = DATA / "test_kaiser_deck.pptx"
    out.write_bytes(buf.getvalue())
    print(f"✓ Wrote {out} ({out.stat().st_size / 1024:.0f} KB)")
